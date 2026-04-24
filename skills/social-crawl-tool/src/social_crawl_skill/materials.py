from __future__ import annotations

import json
from pathlib import Path

import fitz
from docx import Document
from pptx import Presentation


TEXT_EXTENSIONS = {".md", ".txt", ".pdf", ".docx", ".pptx"}
IMAGE_EXTENSIONS = {".jpg", ".jpeg", ".png", ".webp"}
VIDEO_EXTENSIONS = {".mp4", ".mov", ".avi", ".mkv"}


def _read_text_file(path: Path) -> str:
    return path.read_text(encoding="utf-8", errors="ignore")


def _read_pdf(path: Path) -> str:
    doc = fitz.open(path)
    chunks = [page.get_text("text") for page in doc]
    return "\n".join(chunks)


def _read_docx(path: Path) -> str:
    doc = Document(path)
    return "\n".join(p.text for p in doc.paragraphs if p.text.strip())


def _read_pptx(path: Path) -> str:
    prs = Presentation(path)
    texts: list[str] = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text:
                texts.append(shape.text)
    return "\n".join(texts)


def extract_text(path: Path) -> str:
    suffix = path.suffix.lower()
    if suffix in {".md", ".txt"}:
        return _read_text_file(path)
    if suffix == ".pdf":
        return _read_pdf(path)
    if suffix == ".docx":
        return _read_docx(path)
    if suffix == ".pptx":
        return _read_pptx(path)
    return ""


class MaterialIndexer:
    def __init__(self, root: Path) -> None:
        self.root = root

    def build_index(self, source_dir: Path, output_path: Path) -> dict:
        files = []
        text_chunks = []
        for path in source_dir.rglob("*"):
            if not path.is_file():
                continue
            suffix = path.suffix.lower()
            category = "binary"
            if suffix in TEXT_EXTENSIONS:
                category = "text"
            elif suffix in IMAGE_EXTENSIONS:
                category = "image"
            elif suffix in VIDEO_EXTENSIONS:
                category = "video"
            text = extract_text(path) if category == "text" else ""
            if text:
                text_chunks.append({"path": str(path), "excerpt": text[:4000]})
            files.append(
                {
                    "path": str(path),
                    "name": path.name,
                    "category": category,
                    "size": path.stat().st_size,
                }
            )
        payload = {
            "source_dir": str(source_dir),
            "file_count": len(files),
            "files": files,
            "text_chunks": text_chunks,
        }
        output_path.parent.mkdir(parents=True, exist_ok=True)
        output_path.write_text(json.dumps(payload, ensure_ascii=False, indent=2), encoding="utf-8")
        return payload
